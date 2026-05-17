from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from fpdf import FPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from axiom.branding import BRAND, hex_to_rgb, normalize_theme
from axiom.document_plan import deterministic_document_blueprint, normalize_document_blueprint


def render_artifacts(
    frame: pd.DataFrame,
    manifesto: dict[str, Any],
    analysis: dict[str, Any],
    analysis_plan: dict[str, Any],
    run_dir: Path,
    title: str,
    logo_path: Path | None = None,
    brand_guideline: str = "",
) -> dict[str, Path]:
    assets_dir = run_dir / "assets"
    logs_dir = run_dir / "sandbox_logs"
    assets_dir.mkdir(parents=True, exist_ok=True)
    logs_dir.mkdir(parents=True, exist_ok=True)

    manifesto_path = run_dir / "data_manifesto.json"
    summary_path = run_dir / "summary_stats.json"
    manifesto_path.write_text(json.dumps(manifesto, indent=2, default=str), encoding="utf-8")
    summary_path.write_text(json.dumps(analysis, indent=2, default=str), encoding="utf-8")

    visualization_specs = analysis_plan.get("recommended_visualizations", [])
    fallback_blueprint = deterministic_document_blueprint(
        title,
        manifesto,
        visualization_specs,
    )
    draft_blueprint = normalize_document_blueprint(
        analysis_plan.get("document_blueprint"),
        fallback_blueprint,
        chart_count=len(visualization_specs),
    )
    theme = normalize_theme(draft_blueprint.get("theme"))
    chart_paths = _create_charts(frame, analysis, assets_dir, visualization_specs, theme)
    document_blueprint = normalize_document_blueprint(
        analysis_plan.get("document_blueprint"),
        fallback_blueprint,
        chart_count=len(chart_paths),
    )
    document_blueprint["theme"] = theme
    blueprint_path = run_dir / "document_blueprint.json"
    blueprint_path.write_text(json.dumps(document_blueprint, indent=2, default=str), encoding="utf-8")

    report_path = _render_pdf(
        title,
        manifesto,
        analysis,
        analysis_plan,
        document_blueprint,
        chart_paths,
        run_dir / "report.pdf",
        logo_path,
    )
    deck_path = _render_pptx(
        title,
        manifesto,
        analysis,
        analysis_plan,
        document_blueprint,
        chart_paths,
        run_dir / "slide_deck.pptx",
        logo_path,
    )
    workbook_path = _render_xlsx(frame, analysis, run_dir / "raw_data_dashboard.xlsx")
    trace_path = _write_trace(logs_dir / "analysis_trace.py")

    return {
        "manifesto": manifesto_path,
        "summary": summary_path,
        "pdf": report_path,
        "pptx": deck_path,
        "xlsx": workbook_path,
        "document_blueprint": blueprint_path,
        "analysis_trace": trace_path,
    }


def _create_charts(
    frame: pd.DataFrame,
    analysis: dict[str, Any],
    assets_dir: Path,
    visualization_specs: list[dict[str, Any]],
    theme: dict[str, str],
) -> list[Path]:
    sns.set_theme(
        style="whitegrid",
        rc={
            "axes.facecolor": theme["surface"],
            "figure.facecolor": theme["background"],
            "axes.edgecolor": BRAND["border_gray"],
            "axes.labelcolor": theme["text"],
            "text.color": theme["text"],
            "xtick.color": theme["muted_text"],
            "ytick.color": theme["muted_text"],
        },
    )
    chart_paths: list[Path] = []

    for index, spec in enumerate(visualization_specs[:6], start=1):
        path = assets_dir / f"{index:02d}_{_slug(spec.get('title') or spec.get('chart_type', 'chart'))}.png"
        if _render_chart_from_spec(frame, spec, path, theme, index - 1):
            chart_paths.append(path)

    return chart_paths


def _render_chart_from_spec(
    frame: pd.DataFrame,
    spec: dict[str, Any],
    path: Path,
    theme: dict[str, Any],
    palette_index: int,
) -> bool:
    chart_type = spec.get("chart_type")
    title = spec.get("title", str(chart_type).title())
    x = spec.get("x")
    y = spec.get("y")
    aggregation = spec.get("aggregation", "sum")
    palette = theme.get("chart_palette") or [BRAND["primary_blue"], BRAND["electric_blue"], BRAND["ai_purple"]]
    color = palette[palette_index % len(palette)]

    try:
        plt.figure(figsize=(8, 4.5))
        if chart_type == "bar" and x in frame.columns and y in frame.columns:
            plotted = _aggregate(frame, x, y, aggregation).head(12)
            sns.barplot(data=plotted, x=y, y=x, orient="h", color=color)
            plt.xlabel(y)
            plt.ylabel(x)
        elif chart_type == "line" and x in frame.columns and y in frame.columns:
            plotted = _aggregate(frame, x, y, aggregation).sort_values(x)
            sns.lineplot(data=plotted, x=x, y=y, color=color, marker="o")
            plt.xlabel(x)
            plt.ylabel(y)
        elif chart_type == "scatter" and x in frame.columns and y in frame.columns:
            sns.scatterplot(data=frame, x=x, y=y, color=color)
        elif chart_type == "histogram" and x in frame.columns:
            sns.histplot(frame[x].dropna(), kde=True, color=color)
            plt.xlabel(x)
        elif chart_type == "heatmap":
            columns = [column for column in spec.get("columns", []) if column in frame.columns]
            if len(columns) < 2:
                return False
            sns.heatmap(frame[columns].corr(numeric_only=True), annot=True, cmap="coolwarm", center=0)
        else:
            return False

        plt.title(title)
        plt.tight_layout()
        plt.savefig(path, dpi=180)
        plt.close()
        return True
    except Exception:
        plt.close()
        return False


def _aggregate(frame: pd.DataFrame, x: str, y: str, aggregation: str) -> pd.DataFrame:
    if aggregation == "mean":
        return frame.groupby(x, dropna=False, as_index=False)[y].mean().sort_values(y, ascending=False)
    if aggregation == "count":
        return frame.groupby(x, dropna=False, as_index=False)[y].count().sort_values(y, ascending=False)
    if aggregation == "none":
        return frame[[x, y]].dropna()
    return frame.groupby(x, dropna=False, as_index=False)[y].sum().sort_values(y, ascending=False)


def _render_pdf(
    title: str,
    manifesto: dict[str, Any],
    analysis: dict[str, Any],
    analysis_plan: dict[str, Any],
    document_blueprint: dict[str, Any],
    chart_paths: list[Path],
    output_path: Path,
    logo_path: Path | None,
) -> Path:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    _pdf_background(pdf)
    _pdf_logo(pdf, logo_path)
    pdf.set_text_color(*hex_to_rgb(BRAND["silver_white"]))
    pdf.set_font("Helvetica", "B", 18)
    pdf.cell(0, 10, "AXIOM", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*hex_to_rgb(BRAND["soft_gray"]))
    pdf.cell(0, 6, BRAND["tagline"], new_x="LMARGIN", new_y="NEXT")

    pdf.ln(8)
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_text_color(*hex_to_rgb(BRAND["silver_white"]))
    pdf.multi_cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")

    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(*hex_to_rgb(BRAND["soft_gray"]))
    pdf.multi_cell(
        0,
        7,
        f"Rows: {manifesto['row_count']} | Columns: {manifesto['column_count']}",
        new_x="LMARGIN",
        new_y="NEXT",
    )

    first_section = True
    for section in document_blueprint["report_sections"]:
        if not first_section:
            pdf.add_page()
            _pdf_background(pdf)
            _pdf_logo(pdf, logo_path)
        first_section = False
        _render_pdf_section(pdf, section, manifesto, analysis, analysis_plan, chart_paths)

    pdf.output(output_path)
    return output_path


def _render_pptx(
    title: str,
    manifesto: dict[str, Any],
    analysis: dict[str, Any],
    analysis_plan: dict[str, Any],
    document_blueprint: dict[str, Any],
    chart_paths: list[Path],
    output_path: Path,
    logo_path: Path | None,
) -> Path:
    presentation = Presentation()
    for slide_spec in document_blueprint["deck_slides"]:
        _render_ppt_slide(
            presentation,
            slide_spec,
            title,
            manifesto,
            analysis,
            analysis_plan,
            document_blueprint,
            chart_paths,
            logo_path,
        )

    presentation.save(output_path)
    return output_path


def _render_xlsx(frame: pd.DataFrame, analysis: dict[str, Any], output_path: Path) -> Path:
    with pd.ExcelWriter(output_path, engine="xlsxwriter", datetime_format="yyyy-mm-dd") as writer:
        frame.to_excel(writer, sheet_name="Clean Data", index=False)
        pd.DataFrame(analysis["numeric_summary"]).T.to_excel(writer, sheet_name="Numeric Summary")

        workbook = writer.book
        data_sheet = writer.sheets["Clean Data"]
        header_format = workbook.add_format(
            {
                "bold": True,
                "font_color": BRAND["silver_white"],
                "bg_color": BRAND["primary_blue"],
                "border": 1,
                "border_color": BRAND["border_gray"],
            }
        )
        for column_index, column in enumerate(frame.columns):
            data_sheet.write(0, column_index, column, header_format)
            data_sheet.set_column(column_index, column_index, min(max(len(column) + 2, 12), 28))

        if analysis["numeric_columns"]:
            first_numeric = analysis["numeric_columns"][0]
            column_index = frame.columns.get_loc(first_numeric)
            data_sheet.conditional_format(
                1,
                column_index,
                len(frame),
                column_index,
                {
                    "type": "3_color_scale",
                    "min_color": BRAND["deep_blue"],
                    "mid_color": BRAND["electric_blue"],
                    "max_color": BRAND["ai_purple"],
                },
            )

    return output_path


def _render_pdf_section(
    pdf: FPDF,
    section: dict[str, Any],
    manifesto: dict[str, Any],
    analysis: dict[str, Any],
    analysis_plan: dict[str, Any],
    chart_paths: list[Path],
) -> None:
    section_type = section["type"]
    _pdf_heading(pdf, section.get("title", section_type.replace("_", " ").title()))

    if section_type == "executive_summary":
        _pdf_bullets(pdf, analysis.get("insight_candidates", [])[:8] or ["No insight candidates generated."])
    elif section_type == "kpi_summary":
        measures = manifesto.get("derived_measures", []) or analysis_plan.get("derived_measures", [])
        _pdf_bullets(
            pdf,
            [f"{measure.get('name')}: {measure.get('formula')} - {measure.get('description')}" for measure in measures]
            or ["No derived KPI measures were created."],
        )
    elif section_type == "chart_story":
        refs = section.get("chart_refs") or list(range(len(chart_paths)))
        for ref in refs:
            if 0 <= ref < len(chart_paths):
                pdf.set_font("Helvetica", "B", 11)
                pdf.set_text_color(*hex_to_rgb(BRAND["soft_cyan"]))
                pdf.cell(0, 7, chart_paths[ref].stem.replace("_", " ").title(), new_x="LMARGIN", new_y="NEXT")
                pdf.image(str(chart_paths[ref]), x=15, w=180)
                pdf.ln(4)
    elif section_type == "data_model":
        relationships = manifesto.get("relationships", [])
        _pdf_bullets(
            pdf,
            [
                f"{item.get('left_table')}.{item.get('left_column')} -> "
                f"{item.get('right_table')}.{item.get('right_column')} ({item.get('relationship_type')})"
                for item in relationships
            ]
            or ["No relationships were inferred."],
        )
    elif section_type == "data_quality":
        warnings = manifesto.get("anomaly_warnings", [])
        sampling = _sampling_lines(manifesto)
        _pdf_bullets(pdf, warnings + sampling or ["No major data quality warnings were detected."])
    elif section_type == "self_healing":
        metadata = analysis.get("self_healing", {})
        _pdf_bullets(
            pdf,
            [
                f"Enabled: {metadata.get('enabled')}",
                f"Fallback used: {metadata.get('used_fallback')}",
                f"Attempt count: {metadata.get('attempt_count')}",
            ],
        )
    elif section_type == "methodology":
        _pdf_bullets(
            pdf,
            [
                "Data was profiled, semantically modeled, analyzed in a generated analyst workspace, and audited.",
                f"Planner source: {analysis_plan.get('planner_source')}",
            ],
        )


def _pdf_heading(pdf: FPDF, text: str) -> None:
    pdf.set_font("Helvetica", "B", 13)
    pdf.set_text_color(*hex_to_rgb(BRAND["electric_blue"]))
    pdf.cell(0, 9, text, new_x="LMARGIN", new_y="NEXT")


def _pdf_bullets(pdf: FPDF, items: list[str]) -> None:
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*hex_to_rgb(BRAND["silver_white"]))
    for item in items:
        pdf.multi_cell(0, 6, f"- {item}", new_x="LMARGIN", new_y="NEXT")


def _render_ppt_slide(
    presentation: Presentation,
    slide_spec: dict[str, Any],
    title: str,
    manifesto: dict[str, Any],
    analysis: dict[str, Any],
    analysis_plan: dict[str, Any],
    document_blueprint: dict[str, Any],
    chart_paths: list[Path],
    logo_path: Path | None,
) -> None:
    theme = normalize_theme(document_blueprint.get("theme"))
    slide_type = slide_spec["type"]
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _style_slide_background(slide, theme)
    _ppt_logo(slide, logo_path)

    if slide_type == "title":
        _ppt_title(slide, title, Inches(0.75), Inches(2.15), Pt(34))
        _ppt_body(slide, [BRAND["tagline"]], Inches(0.78), Inches(3.05), Inches(8.0), Pt(16), theme["secondary_accent"])
        return

    display_title = slide_spec.get("title", slide_type.replace("_", " ").title())
    if slide_type == "chart" and "chart_ref" in slide_spec and slide_spec["chart_ref"] < len(chart_paths):
        display_title = chart_paths[slide_spec["chart_ref"]].stem.replace("_", " ").title()
    _ppt_title(slide, display_title, Inches(0.75), Inches(0.45), Pt(22))

    if slide_type == "chart" and "chart_ref" in slide_spec and slide_spec["chart_ref"] < len(chart_paths):
        slide.shapes.add_picture(str(chart_paths[slide_spec["chart_ref"]]), Inches(0.75), Inches(1.18), width=Inches(8.5))
    elif slide_type == "executive_summary":
        _ppt_body(slide, analysis.get("insight_candidates", [])[:5], Inches(0.85), Inches(1.25), Inches(8.3), Pt(17), theme["text"])
    elif slide_type == "kpi":
        measures = manifesto.get("derived_measures", []) or analysis_plan.get("derived_measures", [])
        _ppt_body(slide, [f"{m.get('name')}: {m.get('formula')}" for m in measures[:5]], Inches(0.85), Inches(1.25), Inches(8.3), Pt(17), theme["text"])
    elif slide_type == "data_model":
        relationships = manifesto.get("relationships", [])
        _ppt_body(
            slide,
            [
                f"{r.get('left_table')}.{r.get('left_column')} -> {r.get('right_table')}.{r.get('right_column')}"
                for r in relationships[:6]
            ],
            Inches(0.85),
            Inches(1.25),
            Inches(8.3),
            Pt(16),
            theme["text"],
        )
    elif slide_type == "data_quality":
        _ppt_body(slide, (manifesto.get("anomaly_warnings", []) + _sampling_lines(manifesto))[:6], Inches(0.85), Inches(1.25), Inches(8.3), Pt(16), theme["text"])
    elif slide_type == "self_healing":
        metadata = analysis.get("self_healing", {})
        _ppt_body(
            slide,
            [
                f"Enabled: {metadata.get('enabled')}",
                f"Fallback used: {metadata.get('used_fallback')}",
                f"Attempt count: {metadata.get('attempt_count')}",
            ],
            Inches(0.85),
            Inches(1.25),
            Inches(8.3),
            Pt(18),
            theme["text"],
        )


def _ppt_title(slide, text: str, left, top, size) -> None:
    box = slide.shapes.add_textbox(left, top, Inches(8.3), Inches(0.6))
    frame = box.text_frame
    frame.text = text
    frame.paragraphs[0].font.size = size
    frame.paragraphs[0].font.bold = True
    frame.paragraphs[0].font.color.rgb = _ppt_color(BRAND["electric_blue"])


def _ppt_body(slide, lines: list[str], left, top, width, size, color: str = BRAND["presentation_text"]) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(4.8))
    frame = box.text_frame
    frame.clear()
    for index, line in enumerate(lines or ["No content generated."]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.font.size = size
        paragraph.font.color.rgb = _ppt_color(color)


def _sampling_lines(manifesto: dict[str, Any]) -> list[str]:
    lines = []
    for name, table in manifesto.get("tables", {}).items():
        metadata = table.get("source_metadata", {})
        if metadata.get("processing_mode") == "chunked_csv":
            lines.append(
                f"{name}: scanned {metadata.get('total_rows')} rows in "
                f"{metadata.get('chunks_processed')} chunks; analyzed sample of {metadata.get('sample_rows')} rows."
            )
    return lines[:6]


def _pdf_background(pdf: FPDF) -> None:
    pdf.set_fill_color(*hex_to_rgb(BRAND["dark_navy"]))
    pdf.rect(0, 0, 210, 297, style="F")
    pdf.set_draw_color(*hex_to_rgb(BRAND["primary_blue"]))
    pdf.set_line_width(0.8)
    pdf.line(15, 24, 195, 24)


def _pdf_logo(pdf: FPDF, logo_path: Path | None) -> None:
    if logo_path and logo_path.exists():
        pdf.image(str(logo_path), x=165, y=8, w=28)


def _style_slide_background(slide, theme: dict[str, str]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _ppt_color(theme["background"])


def _ppt_logo(slide, logo_path: Path | None) -> None:
    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(8.25), Inches(0.25), width=Inches(1.25))


def _ppt_color(hex_color: str) -> RGBColor:
    return RGBColor(*hex_to_rgb(hex_color))


def _write_trace(output_path: Path) -> Path:
    output_path.write_text(
        "\n".join(
            [
                "# Project Axiom MVP analysis trace",
                "# Self-healing analyst attempts are stored in ../analyst_workspace/.",
                "profile_dataset(frame, source_name)",
                "analyze_dataset(cleaned_frame)",
                "render_artifacts(cleaned_frame, manifesto, analysis, run_dir, title)",
                "",
            ]
        ),
        encoding="utf-8",
    )
    return output_path


def _slug(value: object) -> str:
    text = str(value).strip().lower()
    slug = "".join(character if character.isalnum() else "_" for character in text)
    while "__" in slug:
        slug = slug.replace("__", "_")
    return slug.strip("_") or "chart"
